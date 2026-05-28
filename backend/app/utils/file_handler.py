import os


def extract_from_pptx(filepath: str) -> str:
    """Extract text from PPT/PPTX using python-pptx."""
    # NOTE: python-pptx is not in stdlib; ensure it's in requirements.txt
    from pptx import Presentation

    prs = Presentation(filepath)
    slides_text = []

    for i, slide in enumerate(prs.slides):
        slide_content = [f"[Slide {i + 1}]"]

        # Slide titles + shapes (including text boxes)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text
                if txt and txt.strip():
                    slide_content.append(txt.strip())

        # Speaker notes (if any)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text if notes_frame else ""
            if notes and notes.strip():
                slide_content.append(f"[Notes: {notes.strip()}]")

        # Slide numbers as section markers (best-effort)
        slide_content.append(f"[Slide Number: {i + 1}]")
        slides_text.append("\n".join(slide_content))

    return "\n\n".join(slides_text)


def extract_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".pdf":
        from pdfminer.high_level import extract_text as pdf_extract

        return pdf_extract(filepath) or ""

    if ext in [".docx"]:
        from docx import Document

        doc = Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs if p.text and p.text.strip()])

    if ext in [".pptx", ".ppt"]:
        # python-pptx supports PPTX; .ppt is best-effort (often unreadable without conversion)
        return extract_from_pptx(filepath)

    if ext == ".txt":
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    if ext == ".doc":
        # Fallback for legacy Word formats: best-effort via docx if it was actually docx,
        # otherwise raise a more helpful error.
        try:
            # Some files may be in a docx container with .doc extension
            from docx import Document

            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs if p.text and p.text.strip()])
        except Exception as e:
            raise ValueError("DOC parsing fallback failed. Please convert DOC to DOCX for best results.") from e

    raise ValueError("Unsupported file type")

